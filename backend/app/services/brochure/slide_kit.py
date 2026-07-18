"""Low-level python-pptx drawing primitives shared by pptx_generator.py and
one_pager.py — kept separate so neither module reaches into the other's
private helpers.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.brochure.theme import Theme

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def color(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


def fill_background(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(hex_color)


def add_rect(slide, left, top, width, height, hex_color: str, line: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(hex_color)
    if not line:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int,
    color_hex: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_family: str = "Calibri",
    italic: bool = False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_family
    run.font.color.rgb = color(color_hex)
    return box


def add_bullets(slide, left, top, width, height, items: list[str], *, size: int, color_hex: str, font_family: str):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.name = font_family
        run.font.color.rgb = color(color_hex)
        p.space_after = Pt(6)
    return box


def eyebrow_and_heading(slide, theme: Theme, eyebrow: str, heading: str, *, dark: bool):
    text_color = theme.text_on_dark_hex if dark else theme.text_on_light_hex
    add_text(
        slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.5), eyebrow.upper(),
        size=14, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family,
    )
    add_text(
        slide, Inches(0.7), Inches(0.95), Inches(11.5), Inches(1), heading,
        size=32, color_hex=text_color, bold=True, font_family=theme.font_family,
    )


def style_table_header(cell, theme: Theme, text: str) -> None:
    cell.text = text
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(11)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color(theme.dark_bg_hex)
    cell.text_frame.paragraphs[0].font.color.rgb = color(theme.text_on_dark_hex)


def style_table_cell(cell, theme: Theme, text: str) -> None:
    cell.text = text
    cell.text_frame.paragraphs[0].font.size = Pt(10.5)
    cell.text_frame.paragraphs[0].font.color.rgb = color(theme.text_on_light_hex)


# ─────────────────────────────── Diagonal brand motif ───────────────────────────────
#
# The reference "Market Inventory" template's signature is a diagonal red cut —
# a full panel on the cover, a smaller corner wedge on interior section pages.
# Both are drawn as freeform polygons rather than approximated with a rotated
# rectangle, so the slant is exact rather than eyeballed.


def add_freeform_polygon(slide, points: list[tuple[int, int]], hex_color: str):
    """`points` are absolute EMU coordinates (e.g. Inches(x)) for a closed polygon."""
    start_x, start_y = points[0]
    builder = slide.shapes.build_freeform(start_x, start_y, scale=1)
    builder.add_line_segments(points[1:], close=True)
    shape = builder.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(hex_color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_diagonal_cover_panel(slide, theme: Theme, *, top_band_height=Inches(1.0)):
    """Full-height diagonal red panel for the cover slide — wider at the top,
    narrower at the bottom, revealing more of the background photo as it
    descends (matching the source template's cover)."""
    top = top_band_height
    bottom = SLIDE_HEIGHT
    points = [
        (Inches(0), top),
        (Inches(8.6), top),
        (Inches(5.6), bottom),
        (Inches(0), bottom),
    ]
    return add_freeform_polygon(slide, points, theme.accent_hex)


def add_diagonal_corner_accent(slide, theme: Theme, *, size=Inches(2.6)):
    """Small red diagonal wedge in the top-left corner, used on interior
    section pages (Search Profile, Per Region divider, Project Team) to
    carry the same diagonal brand language without a full panel."""
    points = [
        (Inches(0), Inches(0)),
        (size, Inches(0)),
        (Inches(0), size),
    ]
    return add_freeform_polygon(slide, points, theme.accent_hex)


def add_numbered_badge(slide, left, top, number: int, theme: Theme, *, size=Inches(0.55)):
    """Small solid square with a bold white number — the per-listing badge
    overlaid on the bottom-left corner of its photo in the source template."""
    shape = add_rect(slide, left, top, size, size, theme.accent_hex)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = color(theme.text_on_dark_hex)
    return shape


def add_photo_placeholder(slide, left, top, width, height, label: str, theme: Theme):
    """Light-gray placeholder standing in for a real downloaded photo (§10) —
    labelled so it's obviously a placeholder rather than a broken image."""
    box = add_rect(slide, left, top, width, height, "E9E9EC")
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = color(theme.muted_text_hex)
    return box


def add_stat_icon_row(slide, left, top, width, items: list[tuple[str, str, str]], theme: Theme):
    """Three-column row of small glyph + label + value, mirroring the
    Accessibility / Public Transport / Airport row under each unit photo in
    the source template. `items` is (glyph, label, value)."""
    col_width = Emu(int(width) // len(items)) if items else width
    x = left
    for glyph, label, value in items:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, top, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color(theme.muted_text_hex)
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        crun.text = glyph
        crun.font.size = Pt(14)
        crun.font.bold = True
        crun.font.color.rgb = color(theme.light_bg_hex)

        add_text(
            slide, x - Inches(0.5), top + Inches(0.65), col_width, Inches(0.3), label.upper(),
            size=10, color_hex=theme.text_on_light_hex, bold=True, align=PP_ALIGN.CENTER, font_family=theme.font_family,
        )
        add_text(
            slide, x - Inches(0.5), top + Inches(0.95), col_width, Inches(0.4), value,
            size=10, color_hex=theme.muted_text_hex, align=PP_ALIGN.CENTER, font_family=theme.font_family,
        )
        x = Emu(int(x) + int(col_width))
