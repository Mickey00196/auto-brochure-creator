"""§14 PowerPoint Generator — treated as the primary generation target.

Layout matches the real "Market Inventory" template this engine now targets
(diagonal red/white brand motif, single-photo unit fact sheets, region
grouping with real map pages, a closing Project Team page) rather than the
brochure-style title-card/gallery layout this generator started with. Slide
order: Cover → Search Profile → Per Region divider → [region map, then one
fact-sheet slide per unit] for each region → All Properties Overview →
Comparison (§16, this engine's own computed value-add) → Recommendation →
Project Team.

PDF is produced from this same deck via LibreOffice headless conversion
(pdf_export.py) — "a flattened export of the same slides," not a second
renderer (§14).
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from app.models import Building, Client, Proposal, Unit
from app.models.enums import DeliveryCondition, PricingModel
from app.services.brochure.slide_kit import (
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    add_bullets,
    add_diagonal_corner_accent,
    add_diagonal_cover_panel,
    add_freeform_polygon,
    add_numbered_badge,
    add_photo_placeholder,
    add_rect,
    add_stat_icon_row,
    add_text,
    blank_slide,
    color,
    eyebrow_and_heading,
    fill_background,
    style_table_cell,
    style_table_header,
)
from app.services.brochure.theme import DEFAULT_THEME, Theme
from app.services.comparison import ComparisonRow, build_comparison_table
from app.services.maps import build_region_map_url, fetch_static_map_image
from app.services.qa import QASeverity, run_qa_pass

_DELIVERY_LABELS = {
    DeliveryCondition.TURN_KEY: "Full service.",
    DeliveryCondition.SHELL_AND_CORE: "Shell & core.",
    DeliveryCondition.SHELL_AND_CORE_PLUS: "Shell & core plus.",
    DeliveryCondition.MIXED: "Mixed.",
}


def _fmt_rate(value: float | None) -> str:
    return f"€{value:,.0f}/m²/yr" if value is not None else "TBD"


def _first_parking_addon(unit: Unit):
    return next((a for a in (unit.addons or []) if "parking" in a.name.lower()), None)


def _humanize_tokens(tokens: list[str]) -> str:
    """Client-facing copy should never leak internal snake_case tokens like
    "near_public_transport" verbatim (§9/§24: never ship raw data as copy)."""
    return ", ".join(t.replace("_", " ").capitalize() for t in tokens)


# ─────────────────────────────────────────── Slides ───────────────────────────────────────────


def add_cover_slide(prs: Presentation, proposal: Proposal, client: Client, theme: Theme):
    slide = blank_slide(prs)
    fill_background(slide, "D9D9DC")  # stand-in for the aerial/hero photo behind the panel

    top_band_height = Inches(1.0)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, top_band_height, theme.light_bg_hex)
    add_text(
        slide, Inches(0.5), Inches(0.28), Inches(6), Inches(0.5), proposal.prepared_by or "Your Advisory Firm",
        size=16, color_hex=theme.muted_text_hex, bold=True, font_family=theme.font_family,
    )

    add_diagonal_cover_panel(slide, theme, top_band_height=top_band_height)

    add_text(
        slide, Inches(0.7), Inches(1.6), Inches(7), Inches(0.5), proposal.document_type.upper(),
        size=16, color_hex=theme.text_on_dark_hex, bold=True, font_family=theme.font_family,
    )
    add_text(
        slide, Inches(0.7), Inches(2.15), Inches(7.5), Inches(1.1), client.company_name.upper(),
        size=40, color_hex=theme.text_on_dark_hex, bold=True, font_family=theme.font_family,
    )
    add_text(
        slide, Inches(0.7), Inches(3.15), Inches(7), Inches(0.9), proposal.search_area_label or "",
        size=20, color_hex=theme.text_on_dark_hex, bold=True, font_family=theme.font_family,
    )
    add_text(
        slide, Inches(0.7), Inches(4.05), Inches(4), Inches(0.4), proposal.prepared_at.strftime("%B %Y").upper(),
        size=12, color_hex=theme.text_on_dark_hex, bold=True, font_family=theme.font_family,
    )
    return slide


def add_search_profile_slide(prs: Presentation, client: Client, proposal: Proposal, theme: Theme):
    slide = blank_slide(prs)
    fill_background(slide, theme.light_bg_hex)
    add_diagonal_corner_accent(slide, theme)

    add_text(
        slide, Inches(0.9), Inches(1.5), Inches(9.5), Inches(0.8), "SEARCH PROFILE",
        size=32, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family,
    )
    add_text(
        slide, Inches(0.9), Inches(2.3), Inches(8.5), Inches(1.1),
        "Based on the search profile below, this shortlist covers the current supply of buildings "
        "in the region(s) concerned. Locations that don't fully match the profile may still be "
        "included so nothing of potential interest is left out.",
        size=13, color_hex=theme.muted_text_hex, font_family=theme.font_family,
    )

    brief = client.search_brief or {}
    fields = [("Search area", proposal.search_area_label or brief.get("location", "TBD"))]
    if brief.get("size_m2_min") or brief.get("size_m2_max"):
        fields.append(("Size range", f"{brief.get('size_m2_min', '—')}–{brief.get('size_m2_max', '—')} m²"))
    budget = brief.get("budget_eur_per_m2_year_max") or brief.get("budget_eur_per_m2_year")
    if budget:
        fields.append(("Budget", f"up to €{budget:,.0f}/m²/yr"))
    if brief.get("must_haves"):
        fields.append(("Must-haves", _humanize_tokens(brief["must_haves"])))

    y = Inches(3.7)
    for label, value in fields:
        add_rect(slide, Inches(0.9), y, Inches(9.7), Pt(1), "D0D0D3")
        add_text(slide, Inches(0.9), y + Inches(0.12), Inches(2.5), Inches(0.4), label, size=13, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)
        add_text(slide, Inches(3.5), y + Inches(0.12), Inches(7), Inches(0.4), str(value), size=13, color_hex=theme.muted_text_hex, font_family=theme.font_family)
        y = Inches(y.inches + 0.55)
    add_rect(slide, Inches(0.9), y, Inches(9.7), Pt(1), "D0D0D3")
    return slide


def add_region_divider_slide(prs: Presentation, regions: list[str], theme: Theme):
    slide = blank_slide(prs)
    fill_background(slide, theme.light_bg_hex)
    add_photo_placeholder(slide, Inches(0), Inches(0), Inches(6.2), SLIDE_HEIGHT, "Area photography", theme)

    add_text(slide, Inches(6.7), Inches(2.0), Inches(6), Inches(0.4), "PROPERTIES DIVIDED", size=14, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family)
    add_text(slide, Inches(6.7), Inches(2.4), Inches(6), Inches(0.8), "PER REGION", size=34, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)

    y = Inches(3.5)
    for region in regions:
        add_text(slide, Inches(6.7), y, Inches(5.3), Inches(0.4), region.upper(), size=15, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)
        add_text(slide, Inches(11.9), y, Inches(0.4), Inches(0.4), "›", size=18, color_hex=theme.muted_text_hex, bold=True, font_family=theme.font_family)
        y = Inches(y.inches + 0.5)
        add_rect(slide, Inches(6.7), y, Inches(5.6), Pt(1), "D0D0D3")
        y = Inches(y.inches + 0.15)
    return slide


def add_region_map_slide(prs: Presentation, region_name: str, numbered_buildings: list[tuple[int, Building]], theme: Theme):
    slide = blank_slide(prs)
    fill_background(slide, theme.light_bg_hex)

    map_url = build_region_map_url(numbered_buildings)
    image_bytes = fetch_static_map_image(map_url) if map_url else None
    if image_bytes:
        slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(0), Inches(0), height=SLIDE_HEIGHT, width=Inches(7.5))
    else:
        add_photo_placeholder(
            slide, Inches(0), Inches(0), Inches(7.5), SLIDE_HEIGHT,
            "Map unavailable — configure GOOGLE_MAPS_API_KEY to render real tiles", theme,
        )

    add_text(slide, Inches(7.9), Inches(0.5), Inches(5), Inches(0.4), "PROPERTIES IN THIS REGION", size=13, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family)
    add_text(slide, Inches(7.9), Inches(0.9), Inches(5), Inches(0.6), region_name.upper(), size=24, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)

    y = Inches(1.9)
    for index, building in numbered_buildings:
        add_numbered_badge(slide, Inches(7.9), y, index, theme, size=Inches(0.35))
        add_text(slide, Inches(8.4), y - Inches(0.02), Inches(4.4), Inches(0.4), f"{building.address}, {building.city}", size=13, color_hex=theme.text_on_light_hex, font_family=theme.font_family)
        y = Inches(y.inches + 0.5)
        add_rect(slide, Inches(7.9), y, Inches(4.9), Pt(1), "D0D0D3")
        y = Inches(y.inches + 0.1)
    return slide


def add_unit_card_slide(prs: Presentation, index: int, unit: Unit, row: ComparisonRow, theme: Theme):
    """Single-photo fact sheet — CHARACTERISTICS + REMARKS + stat row —
    matching the source template's per-listing page (one page per unit, not
    the title-card/detail/gallery 3-slide block this generator used before).
    """
    slide = blank_slide(prs)
    fill_background(slide, theme.light_bg_hex)

    photo_w, photo_h = Inches(4.9), Inches(4.0)
    photo = add_photo_placeholder(slide, Inches(0.6), Inches(0.5), photo_w, photo_h, unit.building.name, theme)
    add_numbered_badge(slide, Inches(0.6), Inches(0.5) + photo_h - Inches(0.55), index, theme)

    add_text(slide, Inches(0.6), Inches(4.7), Inches(4.9), Inches(0.4), unit.building.city.upper(), size=13, color_hex=theme.muted_text_hex, bold=True, font_family=theme.font_family)
    add_text(slide, Inches(0.6), Inches(5.05), Inches(4.9), Inches(0.55), unit.building.address, size=24, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)
    add_text(slide, Inches(0.6), Inches(5.55), Inches(4.9), Inches(0.4), unit.space_provider or unit.building.building_type or "", size=15, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family)

    icon_items = [
        ("A", "Accessibility", unit.building.accessibility_note or "n/a"),
        ("T", "Transport", _public_transport_summary(unit)),
        ("F", "Airport", unit.building.airport_note or "n/a"),
    ]
    add_stat_icon_row(slide, Inches(0.6), Inches(6.15), Inches(4.9), icon_items, theme)

    right_x = Inches(6.1)
    add_text(slide, right_x, Inches(0.5), Inches(6.6), Inches(0.4), "CHARACTERISTICS", size=14, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family)
    add_rect(slide, right_x, Inches(0.92), Inches(6.6), Pt(1), "D0D0D3")

    # A unit with every optional characteristic set (parking, energy label,
    # BREEAM, year built) plus 5 remarks lines and the footnote can run to
    # ~9 rows — tight enough that this whole block is sized to fit the
    # worst case within the slide's 7.5in height, not just the common case.
    characteristics = _characteristics_lines(unit)
    y = Inches(1.1)
    for label, value in characteristics:
        # Row height accounts for labels too long to fit their column on one
        # line (e.g. "Rental price parking space") so the next row never
        # overlaps a wrapped second line.
        row_height = 0.33 if len(label) <= 20 else 0.58
        add_text(slide, right_x, y, Inches(3.3), Inches(row_height), label, size=12, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)
        add_text(slide, right_x + Inches(3.4), y, Inches(3.2), Inches(row_height), value, size=12, color_hex=theme.muted_text_hex, font_family=theme.font_family)
        y = Inches(y.inches + row_height)

    y = Inches(y.inches + 0.2)
    add_text(slide, right_x, y, Inches(6.6), Inches(0.4), "REMARKS", size=14, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family)
    y = Inches(y.inches + 0.38)
    add_rect(slide, right_x, y, Inches(6.6), Pt(1), "D0D0D3")
    y = Inches(y.inches + 0.12)

    remarks = _remarks_lines(unit)
    add_bullets(slide, right_x, y, Inches(6.6), Inches(2.2), remarks, size=12.5, color_hex=theme.text_on_light_hex, font_family=theme.font_family)

    # Each bullet line is ~0.31in tall at this size/spacing — position the
    # footnote below the actual bullet count rather than a fixed offset, so
    # it never overlaps the last remark line.
    footnote_y = Inches(y.inches + 0.31 * len(remarks) + 0.15)
    period = "month" if unit.pricing_model == PricingModel.PER_DESK_MONTHLY else "year"
    add_text(
        slide, right_x, footnote_y, Inches(6.6), Inches(0.6),
        f"All prices are excluding VAT.\nAll rental prices are per {period}.",
        size=10.5, color_hex=theme.muted_text_hex, italic=True, font_family=theme.font_family,
    )
    return slide


def _public_transport_summary(unit: Unit) -> str:
    nb = unit.building.neighbourhood
    if not nb or not nb.public_transport:
        return "n/a"
    first = nb.public_transport[0]
    return f"{first.get('station', '')} {first.get('walking_time_min', '')} min".strip()


def _characteristics_lines(unit: Unit) -> list[tuple[str, str]]:
    b = unit.building
    lines: list[tuple[str, str]] = [
        ("Subarea", b.submarket or "—"),
        ("Total surface", f"Approx. {b.total_building_area_m2:,.0f} sqm." if b.total_building_area_m2 else "—"),
        ("Available", f"Approx. {unit.available_area_m2:,.0f} sqm. office space"),
        ("Available from", unit.availability or "TBD"),
    ]
    if unit.parking_ratio:
        lines.append(("Parking ratio", unit.parking_ratio))
    parking_addon = _first_parking_addon(unit)
    if parking_addon:
        lines.append(("Rental price parking space", f"EUR {parking_addon.price:,.0f} per place"))

    is_flex = unit.pricing_model == PricingModel.PER_DESK_MONTHLY
    if not is_flex:
        lines.append((
            "Service charges",
            f"€{unit.service_charge_eur_per_m2_year:,.0f}/m²/yr" if unit.service_charge_eur_per_m2_year is not None else "TBD",
        ))
    if b.energy_label:
        lines.append(("Energy rating", b.energy_label))
    if b.breeam_rating:
        lines.append(("Environmental category", f"BREEAM: {b.breeam_rating}"))
    if b.year_built:
        lines.append(("Year of construction", str(b.year_built)))
    return lines


def _remarks_lines(unit: Unit) -> list[str]:
    is_flex = unit.pricing_model == PricingModel.PER_DESK_MONTHLY
    lines = [f"Type: {unit.space_provider or unit.building.building_type or 'Office'}."]
    lines.append(f"Delivery: {_DELIVERY_LABELS.get(unit.delivery_condition, 'Shell & core.')}")

    if is_flex:
        lines.append(f"Size: {unit.desk_count} desks." if unit.desk_count else "Size: TBD.")
        pricing = f"Pricing: EUR {unit.price_per_desk_month_eur:,.0f} per desk per month." if unit.price_per_desk_month_eur is not None else "Pricing: TBD."
        lines.append(pricing)
    else:
        lines.append(f"Size: {unit.available_area_m2:,.0f} sqm.")
        rent_prefix = "from " if unit.rent_price_type == "from" else ""
        pricing = (
            f"Pricing: {rent_prefix}EUR {unit.rent_eur_per_m2_year:,.0f} per sqm per year."
            if unit.rent_eur_per_m2_year is not None
            else "Pricing: TBD."
        )
        lines.append(pricing)

    lines.append(f"Meeting room: {unit.meeting_room_note or 'n/a'}")
    return lines


def add_all_properties_overview_slide(prs: Presentation, numbered_buildings: list[tuple[int, Building]], theme: Theme):
    slide = blank_slide(prs)
    fill_background(slide, theme.light_bg_hex)

    map_url = build_region_map_url(numbered_buildings)
    image_bytes = fetch_static_map_image(map_url) if map_url else None
    if image_bytes:
        slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(0), Inches(0), height=SLIDE_HEIGHT, width=Inches(7.0))
    else:
        add_photo_placeholder(
            slide, Inches(0), Inches(0), Inches(7.0), SLIDE_HEIGHT,
            "Map unavailable — configure GOOGLE_MAPS_API_KEY to render real tiles", theme,
        )

    add_text(slide, Inches(7.4), Inches(0.5), Inches(5.3), Inches(0.4), "ALL PROPERTIES", size=13, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family)
    add_text(slide, Inches(7.4), Inches(0.9), Inches(5.3), Inches(0.6), "OVERVIEW", size=26, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)

    col_spacing = 2.85
    col_width = Inches(2.3)  # < col_spacing - text_x_offset, so text never reaches the next column's badge
    row_height = Inches(0.65)  # tall enough for a wrapped two-line address
    for i, (index, building) in enumerate(numbered_buildings):
        col = i % 2
        row = i // 2
        x = Inches(7.4 + col_spacing * col)
        y = Inches(2.0 + row_height.inches * row)
        add_numbered_badge(slide, x, y, index, theme, size=Inches(0.3))
        add_text(slide, x + Inches(0.45), y - Inches(0.03), col_width, Inches(0.55), f"{building.city}, {building.address}", size=11, color_hex=theme.text_on_light_hex, font_family=theme.font_family)
    return slide


def add_comparison_slide(prs: Presentation, rows: list[ComparisonRow], theme: Theme):
    """§16: computed all-in rate + estimated annual cost for direct-lease
    rows, and total monthly cost for flex/per-desk rows — this engine's own
    value-add over the source template, which shows neither."""
    slide = blank_slide(prs)
    fill_background(slide, theme.light_bg_hex)
    eyebrow_and_heading(slide, theme, "Decision Support", "Comparison", dark=False)

    headers = ["Location", "Floor", "Area (m²)", "Pricing Model", "All-in / Monthly Cost", "Energy"]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.8))
    table = table_shape.table
    for c, header in enumerate(headers):
        style_table_header(table.cell(0, c), theme, header)

    for r, row in enumerate(rows, start=1):
        if row.is_tbd:
            cost_display = "TBD"
        elif row.pricing_model == "per_desk_monthly":
            cost_display = f"€{row.price_per_desk_month_eur:,.0f}/desk/mo (€{row.monthly_total_eur:,.0f} total/mo)"
        else:
            cost_display = f"€{row.all_in_rate_eur_per_m2_year:,.0f}/m²/yr (€{row.estimated_annual_cost_eur:,.0f}/yr)"

        values = [
            row.building_name,
            row.floor or "—",
            f"{row.available_area_m2:,.0f}",
            "Direct lease" if row.pricing_model == "per_sqm_annual" else "Flex / per desk",
            cost_display,
            row.energy_label or "—",
        ]
        for c, val in enumerate(values):
            style_table_cell(table.cell(r, c), theme, val)
    return slide


def add_recommendation_slide(prs: Presentation, rows: list[ComparisonRow], client: Client, theme: Theme):
    slide = blank_slide(prs)
    fill_background(slide, theme.dark_bg_hex)
    ready_rows = [r for r in rows if not r.is_tbd]
    # Prefer a direct-lease recommendation (comparable €/m²/yr basis across the
    # shortlist); fall back to the cheapest flex/per-desk row if none exist.
    top = next((r for r in ready_rows if r.pricing_model == "per_sqm_annual"), ready_rows[0] if ready_rows else None)

    add_text(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.4), "OUR RECOMMENDATION", size=14, color_hex=theme.accent_hex, bold=True, font_family=theme.font_family)
    if top:
        add_text(slide, Inches(0.9), Inches(1.2), Inches(11), Inches(0.9), f"{top.building_name} — {top.floor or ''}", size=32, color_hex=theme.text_on_dark_hex, bold=True, font_family=theme.font_family)
        if top.pricing_model == "per_sqm_annual":
            reasons = [
                f"Lowest all-in rate of the shortlist at €{top.all_in_rate_eur_per_m2_year:,.0f}/m²/yr",
                f"Estimated annual cost of €{top.estimated_annual_cost_eur:,.0f} for {top.available_area_m2:,.0f} m²",
            ]
        else:
            reasons = [
                f"€{top.price_per_desk_month_eur:,.0f} per desk per month across {top.desk_count} desks",
                f"Estimated total cost of €{top.monthly_total_eur:,.0f} per month",
            ]
        if client.search_brief and client.search_brief.get("must_haves"):
            reasons.append(f"Matches search brief must-haves: {_humanize_tokens(client.search_brief['must_haves'])}")
        add_bullets(slide, Inches(0.9), Inches(2.4), Inches(10.5), Inches(3), reasons, size=16, color_hex=theme.text_on_dark_hex, font_family=theme.font_family)
    else:
        add_text(slide, Inches(0.9), Inches(1.2), Inches(11), Inches(1.5), "All shortlisted units are pending final pricing. Recommendation will follow once pricing is confirmed.", size=20, color_hex=theme.text_on_dark_hex, font_family=theme.font_family)
    return slide


def add_project_team_slide(prs: Presentation, proposal: Proposal, theme: Theme):
    slide = blank_slide(prs)
    fill_background(slide, theme.light_bg_hex)
    points = [
        (Inches(0), Inches(0)),
        (Inches(5.2), Inches(0)),
        (Inches(3.2), SLIDE_HEIGHT),
        (Inches(0), SLIDE_HEIGHT),
    ]
    add_freeform_polygon(slide, points, theme.accent_hex)

    add_text(slide, Inches(0.5), Inches(0.7), Inches(3.5), Inches(1.2), "PROJECT\nTEAM", size=30, color_hex=theme.text_on_dark_hex, bold=True, font_family=theme.font_family)
    add_text(slide, Inches(0.5), Inches(2.0), Inches(3.2), Inches(0.4), "For questions please contact:", size=13, color_hex=theme.text_on_dark_hex, font_family=theme.font_family)

    card_w, card_h = Inches(4.2), Inches(1.6)
    x = Inches(4.0)
    for member in (proposal.project_team or []):
        add_rect(slide, x, Inches(0.7), card_w, card_h, theme.light_bg_hex)
        photo_w = Inches(1.1)
        add_photo_placeholder(slide, x, Inches(0.7), photo_w, card_h, "", theme)
        text_width = card_w - photo_w - Inches(0.3)
        add_text(slide, x + photo_w + Inches(0.15), Inches(0.85), text_width, Inches(0.4), member.get("name", ""), size=15, color_hex=theme.text_on_light_hex, bold=True, font_family=theme.font_family)
        add_text(slide, x + photo_w + Inches(0.15), Inches(1.28), text_width, Inches(0.3), member.get("email", ""), size=9.5, color_hex=theme.muted_text_hex, font_family=theme.font_family)
        add_text(slide, x + photo_w + Inches(0.15), Inches(1.72), text_width, Inches(0.3), member.get("phone", ""), size=10.5, color_hex=theme.muted_text_hex, font_family=theme.font_family)
        x = Inches(x.inches + card_w.inches + 0.25)

    footer_y = Inches(4.5)
    footer_lines = [proposal.prepared_by or ""]
    add_text(slide, Inches(0.5), footer_y, Inches(4), Inches(1.5), "\n".join(footer_lines), size=13, color_hex=theme.text_on_dark_hex, bold=True, font_family=theme.font_family)
    return slide


# ─────────────────────────────────────────── Orchestration ───────────────────────────────────────────


def build_pptx(db: Session, proposal: Proposal, *, theme: Theme = DEFAULT_THEME, require_qa_pass: bool = True) -> Presentation:
    qa_report = run_qa_pass(db, proposal)
    if require_qa_pass and not qa_report.is_export_ready:
        blocking = [i.message for i in qa_report.issues if i.severity == QASeverity.BLOCKING]
        raise ValueError(
            "Proposal is not export-ready — resolve or explicitly acknowledge these QA issues first "
            f"(§8/§24): {blocking}"
        )

    units = proposal.selected_units
    sorted_rows = build_comparison_table(units)
    row_by_unit_id = {r.unit_id: r for r in sorted_rows}
    unit_index = {unit.unit_id: i + 1 for i, unit in enumerate(units)}

    regions: dict[str, list[Unit]] = {}
    for unit in units:
        region_name = unit.building.submarket or "Other"
        regions.setdefault(region_name, []).append(unit)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_cover_slide(prs, proposal, proposal.client, theme)
    add_search_profile_slide(prs, proposal.client, proposal, theme)
    add_region_divider_slide(prs, list(regions.keys()), theme)

    for region_name, region_units in regions.items():
        numbered_buildings = [(unit_index[u.unit_id], u.building) for u in region_units]
        add_region_map_slide(prs, region_name, numbered_buildings, theme)
        for unit in region_units:
            add_unit_card_slide(prs, unit_index[unit.unit_id], unit, row_by_unit_id[unit.unit_id], theme)

    all_numbered_buildings = [(unit_index[u.unit_id], u.building) for u in units]
    add_all_properties_overview_slide(prs, all_numbered_buildings, theme)

    add_comparison_slide(prs, sorted_rows, theme)
    add_recommendation_slide(prs, sorted_rows, proposal.client, theme)
    add_project_team_slide(prs, proposal, theme)

    return prs
